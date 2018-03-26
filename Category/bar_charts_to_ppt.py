from pptx import Presentation

#SLD_LAYOUT_TITLE_AND_CONTENT = 1
prs = Presentation('Growth of Category.pptx')
slide = prs.slides.add_slide(prs.slide_layouts[9])
title_placeholder = slide.shapes.title
title_placeholder = 'Hello'
pic_left  = int(prs.slide_width * 0.15)
pic_top   = int(prs.slide_height * 0.1)
pic_width = int(prs.slide_width * 0.7)
placeholder = slide.placeholders[1]
picture = placeholder.insert_picture('Certificates Others.png')

placeholder1 = slide.placeholders[0]
placeholder1.text = 'sdadas'

#pic = slide.shapes.add_picture('Certificates Others.png', pic_left, pic_top, pic_width) 
prs.save('Growth of Category.pptx')